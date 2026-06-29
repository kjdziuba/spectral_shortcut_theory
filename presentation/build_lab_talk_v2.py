"""Lab-talk deck v2 — image-first, casual, concrete examples.

Builds presentation/spectral_shortcut_lab_talk_v2.pptx.

Design rules:
  - One idea per slide
  - Image takes ≥60% of the slide
  - Title at top, optional one-line caption at bottom
  - Strict grid: every shape on integer or half-integer inches
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path("/home/u37314kd/Projects/spectral_shortcut_theory")
IMG = ROOT / "presentation" / "img"
RESULTS = ROOT / "results"
OUT = ROOT / "presentation" / "spectral_shortcut_lab_talk_v2.pptx"

NAVY = RGBColor(0x1A, 0x36, 0x5D)
BLUE = RGBColor(0x3D, 0x7D, 0xC7)
ORANGE = RGBColor(0xF5, 0x8A, 0x1F)
GREEN = RGBColor(0x4D, 0xB6, 0xAC)
RED = RGBColor(0xE5, 0x73, 0x73)
GREY = RGBColor(0x55, 0x55, 0x55)
DARK = RGBColor(0x22, 0x22, 0x22)
LIGHT = RGBColor(0xF5, 0xF5, 0xF5)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
SW = prs.slide_width; SH = prs.slide_height

# Grid: usable area
USABLE_L = Inches(0.5)
USABLE_R = SW - Inches(0.5)
USABLE_W = USABLE_R - USABLE_L  # 12.33"

TITLE_Y = Inches(0.45)
TITLE_H = Inches(0.7)
CAPTION_Y = Inches(6.7)
CAPTION_H = Inches(0.5)
BODY_TOP = Inches(1.3)
BODY_BOTTOM = Inches(6.5)
BODY_H = BODY_BOTTOM - BODY_TOP


def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color
    return tb


def add_title(slide, text, color=NAVY):
    add_text(slide, USABLE_L, TITLE_Y, USABLE_W, TITLE_H,
             text, size=30, bold=True, color=color,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


def add_caption(slide, text, color=GREY):
    add_text(slide, USABLE_L, CAPTION_Y, USABLE_W, CAPTION_H,
             text, size=14, color=color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_image_centered(slide, path, max_w_in=None, max_h_in=None,
                       y_in=None):
    """Add image centered horizontally. Auto-scale to fit max_w/max_h."""
    pic = slide.shapes.add_picture(str(path), 0, 0)
    w = pic.width; h = pic.height
    if max_w_in is not None and w > Inches(max_w_in):
        s = Inches(max_w_in) / w
        w = int(w * s); h = int(h * s)
    if max_h_in is not None and h > Inches(max_h_in):
        s = Inches(max_h_in) / h
        w = int(w * s); h = int(h * s)
    pic.width = w; pic.height = h
    pic.left = int((SW - w) / 2)
    if y_in is not None:
        pic.top = Inches(y_in)
    else:
        pic.top = int(BODY_TOP + (BODY_H - h) / 2)
    return pic


def add_pagenum(slide, n, total):
    add_text(slide, SW - Inches(1.0), SH - Inches(0.4),
             Inches(0.7), Inches(0.3),
             f"{n} / {total}", size=10, color=GREY,
             align=PP_ALIGN.RIGHT)


N_SLIDES = 7


# ============================================================ #
# Slide 1 — Title (super clean)
# ============================================================ #
s = blank()

# big centered title
add_text(s, USABLE_L, Inches(2.5), USABLE_W, Inches(1.2),
         "The Spectral Shortcut",
         size=68, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, USABLE_L, Inches(3.7), USABLE_W, Inches(0.6),
         "why locking part of a model makes it better",
         size=24, color=GREY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# subtle ribbon decoration
ribbon = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(5.16), Inches(4.7),
                             Inches(3.0), Inches(0.08))
ribbon.line.fill.background()
ribbon.fill.solid(); ribbon.fill.fore_color.rgb = ORANGE

add_text(s, USABLE_L, Inches(5.1), USABLE_W, Inches(0.4),
         "Krzysztof Dziuba   ·   lab meeting",
         size=14, color=GREY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================ #
# Slide 2 — The paradox
# ============================================================ #
s = blank()
add_title(s, "the puzzle")
add_image_centered(s, IMG / "paradox.png", max_w_in=12.0, max_h_in=5.2)
add_caption(s, "same model — totally different results")
add_pagenum(s, 2, N_SLIDES)


# ============================================================ #
# Slide 3 — The setup (pipeline)
# ============================================================ #
s = blank()
add_title(s, "what's inside")
add_image_centered(s, IMG / "pipeline.png", max_w_in=12.5, max_h_in=5.0)
add_caption(s, "a small \"squeezer\" feeds a huge \"brain\" — and the size gap is the whole story")
add_pagenum(s, 3, N_SLIDES)


# ============================================================ #
# Slide 4 — Chain rule + relay metaphor
# ============================================================ #
s = blank()
add_title(s, "gradients move backward — like a relay")

# The equation centered just below the title
add_text(s, USABLE_L, Inches(1.35), USABLE_W, Inches(0.8),
         "∂L/∂θ   =   ∂L/∂ŷ   ·   ∂ŷ/∂Z   ·   ∂Z/∂θ",
         size=32, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE, font="Cambria Math")

# Relay illustration filling the rest
add_image_centered(s, IMG / "relay.png",
                   max_w_in=11.5, max_h_in=4.2, y_in=2.4)

add_caption(s, "three runners pass the same baton: \"residual\" → \"Jacobian\" → \"input\" → θ")
add_pagenum(s, 4, N_SLIDES)


# ============================================================ #
# Slide 5 — What goes wrong (dropped baton)
# ============================================================ #
s = blank()
add_title(s, "and then the first runner drops it", color=RED)
add_image_centered(s, IMG / "relay_drop.png",
                   max_w_in=11.5, max_h_in=5.0)
add_caption(s, "cross-entropy saturates fast → residual = 0 → no gradient ever reaches θ")
add_pagenum(s, 5, N_SLIDES)


# ============================================================ #
# Slide 6 — The shortcut metaphor (two-moons)
# ============================================================ #
s = blank()
add_title(s, "the \"easy line\" beats the \"right curve\"")

add_image_centered(s, IMG / "moons.png",
                   max_w_in=9.0, max_h_in=5.0)
add_caption(s,
            "imagine cancer (blue) vs normal (orange) — gradient descent picks the easy line, "
            "ignores the curve.   In our case the \"easy line\" is morphology, the \"right curve\" is chemistry.")
add_pagenum(s, 6, N_SLIDES)


# ============================================================ #
# Slide 7 — Results + prescription combined
# ============================================================ #
s = blank()
add_title(s, "and the data shows it")

# Two-column layout: result plot + prescription
col_w = Inches(6.0)
col_y = Inches(1.4)
col_h = Inches(4.6)

# LEFT: paired eigenvalues
left_x = Inches(0.4)
pic = s.shapes.add_picture(str(RESULTS / "exp1_1_paired_eigenvalues.png"),
                            left_x, col_y, height=col_h)
# scale to height
pic.width = int(pic.width)  # keep aspect

# Caption under left
add_text(s, left_x, Inches(6.05), Inches(6.5), Inches(0.4),
         "spatial curvature (blue) grows 75×.\n"
         "spectral curvature (red) stays flat.",
         size=12, color=GREY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

# RIGHT: prescription illustration
right_x = SW - Inches(6.4)
add_image_centered_inside = s.shapes.add_picture(
    str(IMG / "prescription.png"),
    right_x, col_y + Inches(0.5),
    width=Inches(6.0))

add_text(s, right_x, Inches(6.05), Inches(6.0), Inches(0.4),
         "freeze.  anchor.  watch.",
         size=14, bold=True, color=DARK,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

# Footer one-liner
add_text(s, USABLE_L, Inches(6.8), USABLE_W, Inches(0.4),
         "don't train everything at once — the math says you can't.",
         size=16, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_pagenum(s, 7, N_SLIDES)


# ============================================================ #
prs.save(str(OUT))
print(f"wrote {OUT}")
print(f"size: {OUT.stat().st_size / 1024:.1f} KB")
