"""Build a short lab-meeting presentation explaining the project.

Audience: spectroscopy folks, math newbies. Goal: convey the WHY,
the chain-rule intuition, the curved-vs-linear shortcut analogy, and
some preliminary results without overwhelming. Re-run this script to
regenerate the PPTX.

Output: presentation/spectral_shortcut_lab_talk.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path("/home/u37314kd/Projects/spectral_shortcut_theory")
RESULTS = ROOT / "results"
OUT = ROOT / "presentation" / "spectral_shortcut_lab_talk.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)


# ----- colors -----
NAVY = RGBColor(0x1A, 0x36, 0x5D)
BLUE = RGBColor(0x2E, 0x5E, 0x8C)
ORANGE = RGBColor(0xD9, 0x6A, 0x29)
GREEN = RGBColor(0x3E, 0x8A, 0x4E)
RED = RGBColor(0xB7, 0x2C, 0x2C)
LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
GREY = RGBColor(0x55, 0x55, 0x55)
DARK = RGBColor(0x22, 0x22, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height


# ============================================================ #
# helpers
# ============================================================ #
def blank(layout_idx=6):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, font="Calibri"):
    """Add a single-paragraph textbox. Returns the textbox."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=18, color=DARK, bullet="•"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"{bullet}  {item}"
        r.font.size = Pt(size)
        r.font.name = "Calibri"
        r.font.color.rgb = color
        p.space_after = Pt(6)
    return tb


def add_image(slide, path, x, y, w=None, h=None):
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def add_banner(slide, title, accent=NAVY):
    """Top color banner with title text."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.95))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    r.font.name = "Calibri"


def add_footer(slide, text="Spectral Shortcut Theory — lab meeting"):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(12), Inches(0.3),
             text, size=10, color=GREY)


def add_box(slide, x, y, w, h, fill=LIGHT, edge=BLUE):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = edge
    box.line.width = Pt(1.5)
    box.text_frame.word_wrap = True
    return box


# ============================================================ #
# Slide 1 — Title
# ============================================================ #
s = blank()
# big colored half-block on the left
left_block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.5), SH)
left_block.line.fill.background()
left_block.fill.solid()
left_block.fill.fore_color.rgb = NAVY

add_text(s, Inches(0.5), Inches(1.2), Inches(3.7), Inches(1.5),
         "spectral\nshortcut\ntheory", size=44, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF))

add_text(s, Inches(0.5), Inches(4.0), Inches(3.7), Inches(0.6),
         "lab meeting", size=14, color=RGBColor(0xE0, 0xE8, 0xF0))

add_text(s, Inches(0.5), Inches(6.7), Inches(3.7), Inches(0.4),
         "Krzysztof Dziuba", size=14, color=RGBColor(0xE0, 0xE8, 0xF0))

add_text(s, Inches(5.0), Inches(1.5), Inches(8), Inches(1.0),
         "Why end-to-end training fails",
         size=36, bold=True, color=NAVY)
add_text(s, Inches(5.0), Inches(2.4), Inches(8), Inches(0.7),
         "in spectral-spatial deep learning",
         size=28, color=DARK)

add_text(s, Inches(5.0), Inches(3.6), Inches(7.8), Inches(2.0),
         "A mathematical diagnosis of why the same network does better "
         "when we DON'T train part of it — and what this means for "
         "FTIR / QCL classification.",
         size=18, color=GREY)

add_text(s, Inches(5.0), Inches(6.6), Inches(8), Inches(0.5),
         "(Joint paper #2; companion to the empirical results)",
         size=12, color=GREY)


# ============================================================ #
# Slide 2 — The empirical paradox
# ============================================================ #
s = blank()
add_banner(s, "The empirical paradox we keep running into")

add_text(s, Inches(0.5), Inches(1.2), Inches(8), Inches(0.5),
         "Same architecture. Same data. Same training budget. Different test F1.",
         size=20, bold=True, color=NAVY)

# Big table-like layout
table_x = Inches(0.7); table_y = Inches(2.1)
row_h = Inches(0.55)
col_label = Inches(5.5); col_val = Inches(1.8); col_note = Inches(4.5)

variants = [
    ("End-to-end learned linear  (joint)", "0.675", "the standard approach", RED),
    ("Frozen random spectral", "0.70", "even random beats joint",   ORANGE),
    ("Frozen PCA-128 spectral", "0.78", "fixed but informative",     ORANGE),
    ("Frozen pretrained spectral (Peak)",    "0.84", "freeze + chemistry-aware", GREEN),
    ("Frozen pretrained spectral (MLP)",     "0.90", "freeze + learned",         GREEN),
    ("Frozen pretrained spectral (Slidewin)","0.95", "best result we have",      GREEN),
    ("Fine-tuned pretrained",                "0.79", "UNfreezing destroys it",   RED),
]
add_text(s, table_x, Inches(1.6), col_label, Inches(0.4),
         "Variant", size=14, bold=True, color=GREY)
add_text(s, table_x + col_label, Inches(1.6), col_val, Inches(0.4),
         "Test F1", size=14, bold=True, color=GREY, align=PP_ALIGN.CENTER)
add_text(s, table_x + col_label + col_val, Inches(1.6), col_note, Inches(0.4),
         "Note", size=14, bold=True, color=GREY)

for i, (name, score, note, color) in enumerate(variants):
    y = table_y + row_h * i
    add_text(s, table_x, y, col_label, row_h, name, size=15, color=DARK)
    add_text(s, table_x + col_label, y, col_val, row_h, score,
             size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s, table_x + col_label + col_val, y, col_note, row_h, note,
             size=13, color=GREY)

# Right callout
box = add_box(s, Inches(8.5), Inches(2.0), Inches(4.4), Inches(4.0))
add_text(s, Inches(8.8), Inches(2.15), Inches(3.9), Inches(0.5),
         "What this is telling us", size=18, bold=True, color=NAVY)
add_bullets(s, Inches(8.8), Inches(2.7), Inches(3.9), Inches(3.2),
            [
                "Joint training underperforms",
                "Random frozen beats joint by 3 points",
                "Pretrained frozen wins by 30 points",
                "Fine-tuning UNDOES the pretraining",
                "Symptom: training procedure failing",
                "Not the data, not the architecture",
            ],
            size=14, color=DARK)

add_footer(s)


# ============================================================ #
# Slide 3 — Setup: the compositional model
# ============================================================ #
s = blank()
add_banner(s, "The setup — a compositional model")

add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
         "Almost every spectral-spatial deep model has this shape:",
         size=18, color=DARK)

# Pipeline boxes
y_pipe = Inches(2.2); h_pipe = Inches(1.4)
box_w = Inches(2.6); gap = Inches(0.7)
x0 = Inches(0.6)

# Input
b1 = add_box(s, x0, y_pipe, box_w, h_pipe, fill=LIGHT, edge=GREY)
add_text(s, x0 + Inches(0.1), y_pipe + Inches(0.1), box_w, Inches(0.5),
         "input", size=14, bold=True, color=GREY)
add_text(s, x0 + Inches(0.1), y_pipe + Inches(0.5), box_w, Inches(0.8),
         "X ∈ R^(H × W × S)\nfor each pixel: a spectrum",
         size=14, color=DARK)

# Spectral
x1 = x0 + box_w + gap
b2 = add_box(s, x1, y_pipe, box_w, h_pipe, fill=LIGHT, edge=BLUE)
add_text(s, x1 + Inches(0.1), y_pipe + Inches(0.1), box_w, Inches(0.5),
         "f_θ  spectral reduction", size=14, bold=True, color=BLUE)
add_text(s, x1 + Inches(0.1), y_pipe + Inches(0.5), box_w, Inches(0.8),
         "per-pixel  R^S → R^K\nlinear, small (C_f params)",
         size=14, color=DARK)

# Spatial
x2 = x1 + box_w + gap
b3 = add_box(s, x2, y_pipe, box_w, h_pipe, fill=LIGHT, edge=ORANGE)
add_text(s, x2 + Inches(0.1), y_pipe + Inches(0.1), box_w, Inches(0.5),
         "g_φ  spatial model", size=14, bold=True, color=ORANGE)
add_text(s, x2 + Inches(0.1), y_pipe + Inches(0.5), box_w, Inches(0.8),
         "CNN / ViT / U-Net\nbig (C_g  ≫  C_f)",
         size=14, color=DARK)

# Output
x3 = x2 + box_w + gap
b4 = add_box(s, x3, y_pipe, box_w, h_pipe, fill=LIGHT, edge=GREEN)
add_text(s, x3 + Inches(0.1), y_pipe + Inches(0.1), box_w, Inches(0.5),
         "ŷ  prediction", size=14, bold=True, color=GREEN)
add_text(s, x3 + Inches(0.1), y_pipe + Inches(0.5), box_w, Inches(0.8),
         "per-pixel class\nlogits then softmax",
         size=14, color=DARK)

# Arrows
for ax in [x0 + box_w, x1 + box_w, x2 + box_w]:
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax + Inches(0.05),
                               y_pipe + Inches(0.55), gap - Inches(0.1),
                               Inches(0.3))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = NAVY
    arrow.line.fill.background()

# Below: capacity asymmetry
add_text(s, Inches(0.5), Inches(4.3), Inches(12), Inches(0.4),
         "The key structural fact: capacity asymmetry",
         size=22, bold=True, color=NAVY)

add_text(s, Inches(0.7), Inches(4.9), Inches(12), Inches(0.6),
         "C_f  ≈  40,000  (linear, ~942 → 128)        "
         "C_g  ≈  13,000,000  (ViT + decoder)        "
         "ratio  C_g / C_f  ≈  325 ×",
         size=18, color=DARK)

# Bottom: the question
box = add_box(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.1),
              fill=RGBColor(0xFF, 0xF8, 0xE8), edge=ORANGE)
add_text(s, Inches(0.8), Inches(5.95), Inches(11.8), Inches(0.4),
         "The question we are answering:", size=16, bold=True, color=ORANGE)
add_text(s, Inches(0.8), Inches(6.35), Inches(11.8), Inches(0.6),
         "why does training f_θ together with g_φ (\"end-to-end\") "
         "fail in a way that freezing f_θ fixes?",
         size=18, color=DARK)

add_footer(s)


# ============================================================ #
# Slide 4 — Chain rule intuition
# ============================================================ #
s = blank()
add_banner(s, "The math intuition — the chain rule decomposition")

add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
         "Gradient flowing back to the spectral parameters factorizes into three pieces:",
         size=18, color=DARK)

# Big equation in the center
add_text(s, Inches(1.5), Inches(2.0), Inches(11), Inches(1.2),
         "∂L / ∂θ   =   ∂L/∂ŷ   ·   ∂ŷ/∂Z   ·   ∂Z/∂θ",
         size=34, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
         font="Cambria Math")

# Labels under each term
y_lab = Inches(3.3); h_lab = Inches(0.45)
add_text(s, Inches(3.5), y_lab, Inches(2.5), h_lab,
         "residual",
         size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_text(s, Inches(6.5), y_lab, Inches(2.5), h_lab,
         "spatial Jacobian",
         size=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(s, Inches(9.5), y_lab, Inches(2.5), h_lab,
         "input",
         size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Three columns explaining the terms
col_w = Inches(4.0); col_y = Inches(4.2); col_h = Inches(2.5)
gap_c = Inches(0.3); col_x0 = Inches(0.4)

# Column 1 — residual
b = add_box(s, col_x0, col_y, col_w, col_h, fill=RGBColor(0xFF, 0xEE, 0xEE),
            edge=RED)
add_text(s, col_x0 + Inches(0.2), col_y + Inches(0.1),
         col_w, Inches(0.4), "1.  ∂L/∂ŷ  —  the error", size=16, bold=True,
         color=RED)
add_bullets(s, col_x0 + Inches(0.2), col_y + Inches(0.55),
            col_w, col_h,
            [
                "how wrong predictions are",
                "softmax(ŷ) − y",
                "when predictions get confident,\nthis goes to 0 EXPONENTIALLY",
                "if this is 0, the whole product is 0",
            ],
            size=13, color=DARK)

# Column 2 — Jacobian
col_x1 = col_x0 + col_w + gap_c
b = add_box(s, col_x1, col_y, col_w, col_h, fill=RGBColor(0xFF, 0xF4, 0xE6),
            edge=ORANGE)
add_text(s, col_x1 + Inches(0.2), col_y + Inches(0.1),
         col_w, Inches(0.4), "2.  ∂ŷ/∂Z  —  what g_φ wants",
         size=16, bold=True, color=ORANGE)
add_bullets(s, col_x1 + Inches(0.2), col_y + Inches(0.55),
            col_w, col_h,
            [
                "Jacobian of the spatial model",
                "tells f_θ \"give me features\nlooking like this\"",
                "depends on g_φ's current state",
                "CHANGES every training step\n→ moving target",
            ],
            size=13, color=DARK)

# Column 3 — input
col_x2 = col_x1 + col_w + gap_c
b = add_box(s, col_x2, col_y, col_w, col_h, fill=RGBColor(0xEC, 0xF7, 0xEC),
            edge=GREEN)
add_text(s, col_x2 + Inches(0.2), col_y + Inches(0.1),
         col_w, Inches(0.4), "3.  ∂Z/∂θ  —  the input",
         size=16, bold=True, color=GREEN)
add_bullets(s, col_x2 + Inches(0.2), col_y + Inches(0.55),
            col_w, col_h,
            [
                "Z = Wᵀ X, so this is just X",
                "the input spectrum",
                "doesn't change during training",
                "delivers signal at full strength",
            ],
            size=13, color=DARK)

add_text(s, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
         "Key insight:  the gradient to the spectral module is killed at term 1.  "
         "Doesn't matter how good the other two are.",
         size=15, bold=True, color=NAVY)

add_footer(s)


# ============================================================ #
# Slide 5 — The shortcut analogy
# ============================================================ #
s = blank()
add_banner(s, "The shortcut — easy answer crowds out the right one")

add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(0.6),
         "From Pezeshki et al. (NeurIPS 2021):  gradient descent picks "
         "the EASIEST separating boundary, not the most useful one.",
         size=16, color=DARK)

# Two-column layout: "Pezeshki idea" vs "our setting"
col_w = Inches(6.1)
col_y = Inches(2.1); col_h = Inches(4.4)

# Left column
b = add_box(s, Inches(0.4), col_y, col_w, col_h,
            fill=RGBColor(0xEE, 0xF1, 0xF7), edge=BLUE)
add_text(s, Inches(0.6), col_y + Inches(0.15), col_w, Inches(0.4),
         "Pezeshki's two-moons", size=18, bold=True, color=BLUE)
add_bullets(s, Inches(0.7), col_y + Inches(0.65), col_w, col_h - Inches(0.7),
            [
                "Two curved classes, almost linearly separable",
                "GD finds the LINEAR boundary",
                "Never finds the curved one — even though\n   curved is more robust on held-out data",
                "Mechanism:  linear feature has bigger\n   eigenvalue in the NTK; the residual decays\n   along it first",
                "Weaker (curved) features get NO\n   gradient signal after that — they starve",
            ],
            size=14, color=DARK)

# Right column
b = add_box(s, Inches(6.85), col_y, col_w, col_h,
            fill=RGBColor(0xFF, 0xF4, 0xE6), edge=ORANGE)
add_text(s, Inches(7.05), col_y + Inches(0.15), col_w, Inches(0.4),
         "Our spectroscopy setting", size=18, bold=True, color=ORANGE)
add_bullets(s, Inches(7.15), col_y + Inches(0.65), col_w, col_h - Inches(0.7),
            [
                "\"Curved\" = chemistry — subtle peak\n   ratios, amide bands, lipid signatures",
                "\"Linear\" = morphology — region size,\n   tissue boundaries, density contrast",
                "Joint training picks MORPHOLOGY because\n   it's the bigger-eigenvalue feature in\n   the high-capacity spatial module",
                "Chemistry never gets the gradient it\n   needs to specialize → spectral shortcut",
                "Same mechanism, different domain",
            ],
            size=14, color=DARK)

# Bottom callout
box = add_box(s, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.6),
              fill=RGBColor(0xFF, 0xEE, 0xEE), edge=RED)
add_text(s, Inches(0.8), Inches(6.75), Inches(12), Inches(0.5),
         "This is why end-to-end training underperforms.  It's a structural property of "
         "compositional models with capacity asymmetry — not a hyperparameter problem.",
         size=14, bold=True, color=RED)

add_footer(s)


# ============================================================ #
# Slide 6 — Preliminary results
# ============================================================ #
s = blank()
add_banner(s, "Preliminary results — what we're seeing")

# Left half: the geometric pathology
add_text(s, Inches(0.4), Inches(1.1), Inches(6), Inches(0.5),
         "1.  The loss landscape is GEOMETRICALLY ill-conditioned",
         size=15, bold=True, color=NAVY)
img = add_image(s, RESULTS / "exp1_1_paired_eigenvalues.png",
                Inches(0.3), Inches(1.6), w=Inches(6.3))

add_text(s, Inches(0.4), Inches(5.6), Inches(6.3), Inches(0.6),
         "Spatial curvature (blue) grows 75× with capacity.  "
         "Spectral curvature (red) stays flat at ~0.02.  "
         "Condition number κ from 40 → 2,400.",
         size=12, color=DARK)

add_text(s, Inches(0.4), Inches(6.4), Inches(6.3), Inches(0.5),
         "SYNTHETIC DATA: per-pixel spectrum (R^64) + smooth class\n"
         "patterns; CNN width D varied across 16 → 8192, 5 seeds.",
         size=10, color=GREY)

# Right half: the equal-info killer test
add_text(s, Inches(6.9), Inches(1.1), Inches(6), Inches(0.5),
         "2.  Joint training collapses BELOW each modality's baseline",
         size=15, bold=True, color=NAVY)
img = add_image(s, RESULTS / "exp1_3v1_summary.png",
                Inches(6.8), Inches(1.6), w=Inches(6.3))

add_text(s, Inches(6.9), Inches(5.6), Inches(6.3), Inches(0.6),
         "Left panel = Bayes-calibrated equal info.  "
         "Joint reaches 0.48 — below frozen (0.61), spectral-only (0.53), "
         "spatial-only (0.53).",
         size=12, color=DARK)

add_text(s, Inches(6.9), Inches(6.4), Inches(6.3), Inches(0.5),
         "SYNTHETIC DATA: spectral signal in u-direction, spatial signal\n"
         "in orthogonal v-direction; α, β calibrated to give equal\n"
         "per-modality baselines.",
         size=10, color=GREY)

add_footer(s)


# ============================================================ #
# Slide 7 — Where this goes
# ============================================================ #
s = blank()
add_banner(s, "Where this goes — the paper and the prescription")

# Two columns
col_w = Inches(6.1)
col_y = Inches(1.4); col_h = Inches(5.3)

# Left: what's in the paper
b = add_box(s, Inches(0.4), col_y, col_w, col_h,
            fill=RGBColor(0xEE, 0xF1, 0xF7), edge=BLUE)
add_text(s, Inches(0.6), col_y + Inches(0.15), col_w, Inches(0.4),
         "What the paper delivers", size=18, bold=True, color=BLUE)
add_bullets(s, Inches(0.7), col_y + Inches(0.65), col_w, col_h - Inches(0.7),
            [
                "Theorem 1: condition number κ(G) ≥ Ω(C_g / C_f) —\n   capacity asymmetry forces ill-conditioning",
                "Theorem 2: under joint SGD the spectral module is\n   GRADIENT-STARVED before it can specialize",
                "EGR diagnostic: ‖∇_θ L‖ / ‖∇_φ L‖ during training —\n   a real-time signal of the pathology",
                "Synthetic verification (Exps 1.1–1.4)",
                "Real-data link (FTIR / QCL companion paper)",
            ],
            size=13, color=DARK)

# Right: prescription
b = add_box(s, Inches(6.85), col_y, col_w, col_h,
            fill=RGBColor(0xEC, 0xF7, 0xEC), edge=GREEN)
add_text(s, Inches(7.05), col_y + Inches(0.15), col_w, Inches(0.4),
         "The prescription",
         size=18, bold=True, color=GREEN)
add_bullets(s, Inches(7.15), col_y + Inches(0.65), col_w, col_h - Inches(0.7),
            [
                "FREEZE the spectral module —\n   removes the timescale gap by construction",
                "Or: PCA / fixed baseline injection —\n   spatial sees stable input even if f_θ drifts",
                "Track EGR during training —\n   when it collapses, training is going wrong",
                "Stop trying to learn everything at once —\n   the math says you can't, on this architecture",
                "Why it matters: every FTIR / QCL / Raman /\n   hyperspectral paper using end-to-end joint\n   training is losing performance on test data",
            ],
            size=13, color=DARK)

# Bottom one-liner
box = add_box(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
              fill=NAVY, edge=NAVY)
add_text(s, Inches(0.7), Inches(7.07), Inches(12), Inches(0.3),
         "One-line takeaway:  in spectral-spatial DL the math says don't "
         "train everything at once — freeze the spectral encoder and the test "
         "F1 goes up.",
         size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

add_footer(s)


# ============================================================ #
prs.save(str(OUT))
print(f"wrote {OUT}")
print(f"size: {OUT.stat().st_size / 1024:.1f} KB")
